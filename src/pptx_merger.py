import re
import json
import logging
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.dml.color import RGBColor

logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

class PPTXMerger:
    JsonPathReg = r"[\$@]|[\$@]\.[0-9A-Za-z_:\-\?\.\[\]\*]+"
    matchRegexpTemplate = re.compile(r".*(\{\{\s*(?P<jsonpath>" + JsonPathReg + r")\s*\}\}).*")
    matchContextControl = re.compile(r"(\{\s*context\s*=\s*(?P<jsonpath>" + JsonPathReg + r")\s*\})")
    _grpShape = re.compile(r"(?s)(\{\s*context\s*=\s*(?P<jsonpath>" + JsonPathReg + r")\s*dir\s*=\s*(?P<dir>\d+)(\s+gap=(?P<gap>\d+))?\s*\})")

    @staticmethod
    def render(data_file, template_file, out_file):
        visited_layouts = []
        logger.debug(f"! Regexp for templates = {PPTXMerger.matchRegexpTemplate}")
        logger.debug(f"! Regexp for controls = {PPTXMerger.matchContextControl}")
        logger.debug(f"! Regexp for group shape controls = {PPTXMerger._grpShape}")

        ppt_template = Presentation(template_file)
        ppt_new = Presentation(template_file)

        # Remove all slides to retain the master slide layouts
        while ppt_new.slides:
            slide = ppt_new.slides[0]
            logger.debug(f"√ Removing slide {PPTXMerger.pretty_print_slide(slide)}")
            ppt_new.slides.remove(slide)

        with open(data_file, 'r') as f:
            json_data = json.load(f)

        for src_slide in ppt_template.slides:
            PPTXMerger.process_slide(template_file, out_file, src_slide, ppt_new, json_data, None, visited_layouts)

        ppt_new.save(out_file)
        return

    @staticmethod
    def pretty_print_slide(slide):
        return f"{slide.name} #{slide.slide_id}"

    @staticmethod
    def pretty_print_shape(shape):
        return shape.name

    @staticmethod
    def process_slide(src_path, dest_path, src_slide, ppt_new, root_json, context_json, visited_layouts):
        new_slide = PPTXTools.create_slide(ppt_new, src_slide, visited_layouts)
        logger.debug(f":: ------ New Slide from src [{src_path}({src_slide.slide_id})] to [{dest_path}({new_slide.slide_id})] ")

        control_data = PPTXMerger.find_control_json_path(new_slide.shapes)
        if control_data:
            if isinstance(control_data, PageControlData):
                logger.debug(f"√ Slide [{PPTXMerger.pretty_print_slide(new_slide)}] has a PageControl - jsonPath context -> {control_data.json_path}")
                new_slide.shapes._spTree.remove(control_data.shape._element)
                result = jsonpath_rw.parse(control_data.json_path).find(root_json)
                for match in result:
                    PPTXMerger.process_slide(src_path, dest_path, new_slide, ppt_new, root_json, match.value, visited_layouts)
                ppt_new.slides.remove(new_slide)
            else:
                logger.error(f"! Slide [{PPTXMerger.pretty_print_slide(new_slide)}] has a control but it is not a PageControlData")
        else:
            PPTXMerger.process_all_shapes(new_slide, root_json, context_json)

    @staticmethod
    def process_all_shapes(slide, root_json, context_json):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                PPTXMerger.process_group_shape(shape)
            elif shape.has_text_frame:
                if PPTXMerger.has_template(shape):
                    logger.debug(f"√ {shape} is a templated shape")
                    PPTXMerger.change_text(shape, root_json, context_json)
                else:
                    logger.debug(f"✖ '{shape.text}' did not match")
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                if PPTXMerger.has_control(shape.table.cell(0, 0)):
                    logger.debug(f"√ we have a table - {shape.table.cell(0, 0).text}")
                    PPTXMerger.iterate_table(shape.table, root_json, context_json)
            else:
                logger.debug(f"✖ {shape} is not TextHolder")

    @staticmethod
    def process_group_shape(group_shape):
        logger.debug(f"⸮ Inspecting GroupShape[{group_shape.name}]...")
        control_data = PPTXMerger.find_control_json_path(group_shape.shapes)
        if control_data:
            if isinstance(control_data, GroupShapeControlData):
                logger.debug(f"√ Found the GroupShape[{group_shape.name}] with control fields")
                new_group_shape = group_shape.parent.shapes.add_group_shape()
                new_group_shape.left = group_shape.left
                new_group_shape.top = group_shape.top
                new_group_shape.width = group_shape.width
                new_group_shape.height = group_shape.height
                # Process shapes within the group
                for shape in group_shape.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        new_shape = new_group_shape.shapes.add_shape(shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height)
                        new_shape.text = shape.text
                    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        PPTXMerger.process_group_shape(shape)
                    elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                        new_shape = new_group_shape.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                        new_shape.text = shape.text
                    else:
                        logger.error(f"The group {PPTXMerger.pretty_print_shape(group_shape)} has a shape {PPTXMerger.pretty_print_shape(shape)} that did not match")
            group_shape._spTree.remove(control_data.shape._element)
        else:
            logger.debug(f"✖ GroupShape[{group_shape.name}] - no control, ignoring")

    @staticmethod
    def iterate_table(table, root_json, context_json):
        first_cell_text = table.cell(0, 0).text
        match = PPTXMerger.matchContextControl.search(first_cell_text)
        if match:
            table_context_json_path = match.group("jsonpath")
            control_string = match.group(0)
            json_node, table_context_json_path = PPTXMerger.node_and_query(table_context_json_path, root_json, context_json)
            result = jsonpath_rw.parse(table_context_json_path).find(json_node)
            for match in result:
                new_row = table.add_row()
                for cell, json_node in zip(new_row.cells, match.value):
                    PPTXMerger.change_text(cell, root_json, json_node)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            table.rows[1]._element.getparent().remove(table.rows[1]._element)
            table.cell(0, 0).text = first_cell_text.replace(control_string, "")

    @staticmethod
    def change_text(shape, root_json, context_json):
        text = shape.text
        match = PPTXMerger.matchRegexpTemplate.search(text)
        if match:
            replacing_text = match.group(1)
            json_query = match.group("jsonpath")
            logger.debug(f"√ found = [{replacing_text}] jsonpath = [{json_query}]")
            json_node, json_path = PPTXMerger.node_and_query(json_query, root_json, context_json)
            result = jsonpath_rw.parse(json_path).find(json_node)
            for match in result:
                new_text = text.replace(replacing_text, match.value)
                logger.debug(f"√ dataText = [{match.value}] newText = [{new_text}]")
                shape.text = new_text

    @staticmethod
    def node_and_query(json_query, root_json, context_json):
        if json_query.startswith('$'):
            return root_json, json_query
        elif json_query.startswith('@'):
            if context_json:
                return context_json, '$' + json_query[1:]
            else:
                new_json_path = '$' + json_query[1:]
                logger.warn(f"! jsonPath starts is {json_query} but the context object is empty. Using root object instead (eg: {new_json_path})")
                return root_json, new_json_path

    @staticmethod
    def find_control_json_path(shapes):
        for shape in shapes:
            if shape.has_text_frame:
                text = shape.text
                logger.debug(f"⸮ inspecting - Shape[{shape.name}] `{text}`")
                match = PPTXMerger._grpShape.search(text)
                if match:
                    json_path = match.group("jsonpath")
                    direction = match.group("dir")
                    gap = int(match.group("gap")) if match.group("gap") else 0
                    control_text = match.group(0)
                    logger.debug(f"√ Match - Shape[{shape.name}] control=`{control_text}` jp=`{json_path}` dir={direction} gap={gap}")
                    return GroupShapeControlData(shape, control_text, json_path, int(direction), gap)
                match = PPTXMerger.matchContextControl.search(text)
                if match:
                    json_path = match.group("jsonpath")
                    control_text = match.group(0)
                    logger.debug(f"√ Match - Shape[{shape.name}] control=`{control_text}` jp=`{json_path}`")
                    return PageControlData(shape, control_text, json_path)
                logger.debug(f"✖ shape:{shape.name} {text} did not have a controlData")
        return None

    @staticmethod
    def has_template(shape):
        return PPTXMerger.matchRegexpTemplate.search(shape.text) is not None

    @staticmethod
    def has_control(shape):
        return PPTXMerger.matchContextControl.search(shape.text) is not None


class ControlData:
    def __init__(self, shape, control_text, json_path):
        self.shape = shape
        self.control_text = control_text
        self.json_path = json_path


class PageControlData(ControlData):
    pass


class GroupShapeControlData(ControlData):
    def __init__(self, shape, control_text, json_path, direction, gap):
        super().__init__(shape, control_text, json_path)
        self.direction = direction
        self.gap = gap


class ImageControlData(ControlData):
    pass


class PPTXTools:
    @staticmethod
    def create_slide(prs, src_slide, visited_layouts):
        slide_layout = src_slide.slide_layout
        if slide_layout not in visited_layouts:
            visited_layouts.append(slide_layout)
        slide = prs.slides.add_slide(slide_layout)
        PPTXMerger.copy_slide_content(src_slide, slide)
        return slide

    @staticmethod
    def copy_slide_content(src_slide, dest_slide):
        for shape in src_slide.shapes:
            if shape.has_text_frame:
                new_shape = dest_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                new_shape.text = shape.text
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_stream = shape.image.blob
                dest_slide.shapes.add_picture(image_stream, shape.left, shape.top, shape.width, shape.height)
            # Add more shape types as needed