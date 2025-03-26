from pptx import Presentation
from pptx.util import Inches
import json
import os

class PPTXTools:

    @staticmethod
    def list_slide_layouts(template_path):
        print(f":: Slide Layouts for {os.path.abspath(template_path)}")
        prs = Presentation(template_path)
        for i, master in enumerate(prs.slide_master.slide_layouts):
            print(f"  :: Master [{i} {master.name}]")
            for layout in master.slide_master.slide_layouts:
                print(f"    Name: {layout.name} - Type: {layout}")

    @staticmethod
    def create_slide(prs, src_slide, visited_layouts):
        slide_layout = src_slide.slide_layout
        if slide_layout not in visited_layouts:
            visited_layouts.append(slide_layout)
        slide = prs.slides.add_slide(slide_layout)
        PPTXTools.copy_slide_content(src_slide, slide)
        return slide

    @staticmethod
    def clone_pptx(src_file, dest_file):

        prs = Presentation(src_file)
        prs.save(dest_file)

    @staticmethod
    def clone_ppt_slides(src_file, dest_file):
        prs_src = Presentation(src_file)
        prs_dest = Presentation()
        prs_dest.slide_width = prs_src.slide_width
        prs_dest.slide_height = prs_src.slide_height

        visited_layouts = []

        for slide in prs_src.slides:
            new_slide = prs_dest.slides.add_slide(slide.slide_layout)
            PPTXTools.copy_slide_content(slide, new_slide)

        prs_dest.save(dest_file)

    @staticmethod
    def copy_slide_content(src_slide, dest_slide):
        for shape in src_slide.shapes:
            if shape.has_text_frame:
                new_shape = dest_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                new_shape.text = shape.text
            elif shape.shape_type == 13:  # Picture
                image_stream = shape.image.blob
                dest_slide.shapes.add_picture(image_stream, shape.left, shape.top, shape.width, shape.height)
            # Add more shape types as needed

class JsonYamlTools:

    @staticmethod
    def parse_json(s):
        try:
            return json.loads(s)
        except json.JSONDecodeError as e:
            print(f"Error parsing JSON: {e}")
            return None

    @staticmethod
    def read_file_to_json(data_file):
        with open(data_file, 'r') as file:
            file_contents = file.read()
        return JsonYamlTools.parse_json(file_contents)